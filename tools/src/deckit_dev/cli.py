from __future__ import annotations

import argparse
import json
import re
import shutil
import zipfile
from pathlib import Path

import yaml


REPO_ROOT = Path(__file__).resolve().parents[3]
DEFAULT_PLUGIN = REPO_ROOT / "plugins" / "deckit"
DEFAULT_RUNS_DIR = Path.cwd() / "outputs"
DEFAULT_MARKETPLACE = REPO_ROOT / ".agents" / "plugins" / "marketplace.json"
TEXT_SOURCE_SUFFIXES = {".md", ".markdown", ".txt"}
PRODUCTION_MODES = ("image-first",)
DELIVERY_TARGETS = ("pptx", "pdf")


def _load_skill_frontmatter(skill_md: Path) -> dict[str, object]:
    text = skill_md.read_text(encoding="utf-8")
    if not text.startswith("---\n"):
        raise ValueError(f"{skill_md} is missing YAML frontmatter")
    try:
        _, frontmatter, _ = text.split("---\n", 2)
    except ValueError as exc:
        raise ValueError(f"{skill_md} has malformed YAML frontmatter") from exc
    data = yaml.safe_load(frontmatter)
    if not isinstance(data, dict):
        raise ValueError(f"{skill_md} frontmatter is not a mapping")
    return data


def inspect_plugin(plugin: Path) -> int:
    manifest_path = plugin / ".codex-plugin" / "plugin.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    skills_dir = plugin / "skills"
    skill_files = sorted(skills_dir.glob("*/SKILL.md"))

    print(f"plugin: {manifest.get('name')}")
    print(f"version: {manifest.get('version')}")
    print(f"skills: {len(skill_files)}")

    for skill_md in skill_files:
        data = _load_skill_frontmatter(skill_md)
        name = data.get("name")
        description = data.get("description")
        if not name or not description:
            raise ValueError(f"{skill_md} must include name and description")
        print(f"- {name}: {skill_md.parent.relative_to(plugin)}")

    return 0


def inspect_marketplace(marketplace: Path) -> int:
    marketplace = marketplace.resolve()
    data = json.loads(marketplace.read_text(encoding="utf-8"))
    plugins = data.get("plugins")
    if not isinstance(plugins, list):
        raise ValueError(f"{marketplace} must contain a plugins list")

    print(f"marketplace: {data.get('name')}")
    print(f"plugins: {len(plugins)}")

    root = marketplace.parents[2]
    for plugin in plugins:
        name = plugin.get("name")
        source = plugin.get("source", {})
        path_value = source.get("path")
        if not name or not path_value:
            raise ValueError(f"marketplace plugin entry is missing name or source.path: {plugin}")
        plugin_path = (root / path_value).resolve()
        manifest_path = plugin_path / ".codex-plugin" / "plugin.json"
        if not manifest_path.is_file():
            raise FileNotFoundError(f"plugin manifest not found for {name}: {manifest_path}")
        try:
            display_path = plugin_path.relative_to(root)
        except ValueError:
            display_path = plugin_path
        print(f"- {name}: {display_path}")

    return 0


def _slugify(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    slug = re.sub(r"-{2,}", "-", slug)
    if not slug:
        raise ValueError("run name must contain at least one letter or digit")
    return slug


def _ensure_inside(parent: Path, child: Path) -> None:
    parent = parent.resolve()
    child = child.resolve()
    if child != parent and parent not in child.parents:
        raise ValueError(f"path is outside expected directory: {child}")


def _classify_source(source: str) -> str:
    """Return one of: 'url', 'pdf', 'text'."""
    from deckit_dev.ingest import is_url

    if is_url(source):
        return "url"
    suffix = Path(source).suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix in TEXT_SOURCE_SUFFIXES:
        return "text"
    raise ValueError(
        f"unsupported source: {source}. Supported: text files ({sorted(TEXT_SOURCE_SUFFIXES)}), .pdf files, and http(s) URLs."
    )


def _slug_from_source(source: str, kind: str) -> str:
    if kind == "url":
        from urllib.parse import urlparse

        parsed = urlparse(source)
        host = parsed.netloc.replace(".", "-")
        path = parsed.path.strip("/").replace("/", "-") or "index"
        return _slugify(f"{host}-{path}")
    return _slugify(Path(source).stem)


def new_run(
    source: str,
    name: str | None,
    runs_dir: Path,
    force: bool,
    mode: str | None,
    budget: str | None,
    delivery_target: str | None,
    target_slides: int | None,
) -> int:
    mode = mode or "image-first"
    kind = _classify_source(source)

    if kind == "text":
        source_path = Path(source).resolve()
        if not source_path.is_file():
            raise FileNotFoundError(f"source file does not exist: {source_path}")
    elif kind == "pdf":
        source_path = Path(source).resolve()
        if not source_path.is_file():
            raise FileNotFoundError(f"pdf source does not exist: {source_path}")
    else:
        source_path = None

    run_name = _slugify(name) if name else _slug_from_source(source, kind)
    runs_dir = runs_dir.resolve()
    run_dir = runs_dir / run_name
    _ensure_inside(runs_dir, run_dir)

    if run_dir.exists() and not force:
        raise FileExistsError(f"run already exists: {run_dir}. Use --force to reuse it.")

    artifacts: dict[str, str] = {
        "deck_brief": "work/deck-brief.md",
        "storyboard": "work/storyboard.md",
        "prompt_readme": "prompts/README.md",
        "prompt_files": "prompts/<slide-id>.md",
        "generated_slides": "assets/generated-slides/<slide-id>.png",
    }
    artifacts["dist"] = "dist/"
    artifacts["preview"] = "dist/preview.png"
    artifacts["review"] = "dist/review.md"

    children: list[str] = ["source", "work", "prompts", "assets/generated-slides", "dist"]
    for child in children:
        (run_dir / child).mkdir(parents=True, exist_ok=True)

    if kind == "text":
        target_source = run_dir / "source" / f"input{source_path.suffix.lower()}"
        if target_source.exists() and not force:
            raise FileExistsError(f"source already exists: {target_source}. Use --force to replace it.")
        shutil.copy2(source_path, target_source)
        original_path = str(source_path)
    else:
        target_source = run_dir / "source" / "input.md"
        if target_source.exists() and not force:
            raise FileExistsError(f"source already exists: {target_source}. Use --force to replace it.")
        from deckit_dev.ingest import ingest_pdf, ingest_url

        if kind == "pdf":
            ingest_pdf(source_path, target_source)
            original_path = str(source_path)
        else:
            ingest_url(source, target_source)
            original_path = source

    manifest: dict[str, object] = {
        "name": run_name,
        "source": {
            "original_path": original_path,
            "local_path": str(target_source.relative_to(run_dir)),
            "type": kind,
        },
        "production_mode": mode,
        "budget_mode": budget,
        "requested_output": {
            "delivery_target": delivery_target,
            "target_slide_count": target_slides,
        },
        "artifacts": artifacts,
    }
    (run_dir / "run.json").write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")

    print(f"created run: {run_dir}")
    print(f"source ({kind}): {target_source}")
    print(f"production_mode: {mode}")
    if budget is None:
        print("budget_mode: not set")
    else:
        print(f"budget_mode: {budget}")
    if delivery_target is None:
        print("delivery_target: not set")
    else:
        print(f"delivery_target: {delivery_target}")
    if target_slides is None:
        print("target_slide_count: not set")
    else:
        print(f"target_slide_count: {target_slides}")
    print("next artifacts:")
    for artifact in artifacts.values():
        print(f"- {artifact}")
    return 0


SLIDE_ID_PATTERN = re.compile(r"^\d{2}_[a-z0-9_]+$")
SLIDE_HEADING_PATTERN = re.compile(r"^##\s+(\S+)\s*$", re.MULTILINE)

TITLE_LENGTH_MAX = 80

KNOWN_ARCHETYPES = {
    "cover", "thesis", "timeline", "comparison", "evidence cards",
    "process", "map", "data chart", "tension", "closing",
}

SUPPORT_COUNT_BANDS: dict[str, tuple[int, int] | None] = {
    "cover": None,
    "closing": None,
    "process": (2, 7),
    "evidence cards": (2, 10),
    "timeline": (2, 10),
    "thesis": (2, 6),
}
SUPPORT_COUNT_DEFAULT = (2, 4)

SLIDE_COUNT_BANDS = {
    "quick": (5, 7),
    "balanced": (7, 10),
    "premium": (8, 14),
}

NATIVE_PPTX_SCRIPT_PATTERNS: tuple[tuple[str, str], ...] = (
    (r"^\s*(from|import)\s+pptx\b", "imports python-pptx"),
    (r"\bPresentation\s*\(", "constructs a PowerPoint presentation"),
    (r"\bslide\.shapes\b", "uses native PowerPoint slide shapes"),
    (r"\bshapes\.add_(textbox|chart|table|shape|connector)\b", "adds native PowerPoint objects"),
    (r"\btext_frame\b", "edits native PowerPoint text frames"),
    (r"\bMSO_SHAPE\b|\bXL_CHART_TYPE\b", "uses native PowerPoint shape/chart constants"),
)

FORBIDDEN_PROMPT_PATTERNS: tuple[tuple[str, str], ...] = (
    (r"\bnative[- ]pptx\b", "asks for native-PPTX output"),
    (r"\bpptx[- ]native\b", "asks for PPTX-native output"),
    (r"\beditable\s+(powerpoint|pptx|ppt)\b", "asks for editable PowerPoint output"),
    (r"\b(powerpoint|pptx|ppt)\s+(text boxes|shapes|charts|layouts)\b", "asks for native PowerPoint objects"),
    (r"\bshape[- ]by[- ]shape\b", "asks for shape-by-shape slide construction"),
    (r"\bpython[- ]pptx\b", "mentions python-pptx in a production artifact"),
    (r"\bhybrid\s+mode\b", "asks for disabled hybrid mode"),
    (r"\buse\s+(codex\s+)?presentations\b", "delegates to a native presentation tool"),
)

NATIVE_PPTX_NEGATION_PATTERN = re.compile(
    r"\b(no|not|never|without|forbidden|disabled|unsupported|does not|do not|is not|must not)\b|不支持|不要|不能|禁止|不是",
    re.IGNORECASE,
)


def _scan_brief(brief_path: Path) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    if not brief_path.is_file():
        findings.append(("error", "BRIEF-001", "work/deck-brief.md missing"))
        return findings
    text = brief_path.read_text(encoding="utf-8").lower()
    for token, rule_id, label in (
        ("thesis", "BRIEF-THESIS", "thesis"),
        ("audience", "BRIEF-AUDIENCE", "audience"),
        ("arc", "BRIEF-ARC", "arc"),
    ):
        if not re.search(rf"^##\s+[^\n]*{token}", text, re.MULTILINE):
            findings.append(("warn", rule_id, f"deck-brief.md missing a heading mentioning '{label}'"))
    return findings


def _parse_storyboard(sb_path: Path) -> tuple[list[tuple[str, str]], list[tuple[str, str, str]]]:
    """Return (slide_blocks, findings). Each slide_block is (slide_id, body_text)."""
    findings: list[tuple[str, str, str]] = []
    if not sb_path.is_file():
        findings.append(("error", "SB-001", "work/storyboard.md missing"))
        return [], findings
    text = sb_path.read_text(encoding="utf-8")
    slide_blocks: list[tuple[str, str]] = []
    matches = list(SLIDE_HEADING_PATTERN.finditer(text))
    if not matches:
        findings.append(("error", "SB-002", "no slide sections (## headings) found in storyboard"))
        return [], findings
    for idx, match in enumerate(matches):
        slide_id = match.group(1)
        if slide_id.lower() in {"deck", "meta"} or slide_id.startswith("Deck"):
            continue
        body_start = match.end()
        body_end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        body = text[body_start:body_end]
        if not SLIDE_ID_PATTERN.match(slide_id):
            findings.append(("warn", "SLIDE-ID-FORMAT", f"slide id '{slide_id}' does not match NN_slug format"))
        slide_blocks.append((slide_id, body))
    return slide_blocks, findings


def _extract_field_value(body: str, name: str) -> str:
    pattern = rf"^-\s+\*\*{re.escape(name)}\*\*:\s*(.+)$"
    m = re.search(pattern, body, re.MULTILINE | re.IGNORECASE)
    return m.group(1).strip() if m else ""


def _extract_markdown_field(body: str, names: tuple[str, ...]) -> str:
    """Extract a storyboard bullet field, including optional indented continuation lines."""
    names_pattern = "|".join(re.escape(name) for name in names)
    pattern = rf"^-\s+\*\*(?:{names_pattern})\*\*:\s*(.*)$"
    match = re.search(pattern, body, re.MULTILINE | re.IGNORECASE)
    if not match:
        return ""

    lines = [match.group(1).strip()]
    body_lines = body.splitlines()
    start_line = body[: match.start()].count("\n")
    for line in body_lines[start_line + 1 :]:
        if re.match(r"^-\s+\*\*[^*]+\*\*:", line):
            break
        if line.strip() == "":
            if lines and lines[-1] != "":
                lines.append("")
            continue
        if line.startswith((" ", "\t")):
            lines.append(line.strip())
            continue
        break
    return "\n".join(line for line in lines).strip()


def _extract_support_points(body: str) -> list[str]:
    sp_match = re.search(
        r"\*\*support points?\*\*[^\n]*\n((?:[ \t]*[-*][ \t]+[^\n]*\n?)+)",
        body,
        re.IGNORECASE,
    )
    if not sp_match:
        return []
    return [
        re.sub(r"^[ \t]*[-*][ \t]+", "", line).strip()
        for line in sp_match.group(1).splitlines()
        if re.match(r"^[ \t]*[-*][ \t]+", line)
    ]


def _build_speaker_notes(slide_id: str, body: str) -> str:
    """Build PowerPoint speaker notes from storyboard speaker guidance.

    Prefer explicit storyboard notes when present. For older storyboards that only
    contain Presenter intent, synthesize a concise talk track from existing fields
    so every packaged slide still has usable presenter guidance.
    """
    explicit_notes = _extract_markdown_field(body, ("Speaker notes", "Presenter notes", "Notes"))
    if explicit_notes:
        return explicit_notes

    title = _extract_field_value(body, "Title")
    primary_job = _extract_field_value(body, "Primary job")
    core_claim = _extract_field_value(body, "Core claim")
    presenter_intent = _extract_field_value(body, "Presenter intent")
    support_points = _extract_support_points(body)

    lines: list[str] = []
    if title:
        lines.append(f"Slide: {title}")
    else:
        lines.append(f"Slide: {slide_id}")

    if presenter_intent:
        lines.extend(["", f"Presenter cue: {presenter_intent}"])
    if core_claim:
        lines.extend(["", f"Core message: {core_claim}"])
    if support_points:
        lines.extend(["", "Talk track:"])
        lines.extend(f"- {point}" for point in support_points)
    if primary_job:
        lines.extend(["", f"Close / transition goal: {primary_job}"])

    return "\n".join(lines).strip()


def _normalize_archetype(value: str) -> str:
    """Strip parenthetical clarifications and lowercase. 'Thesis (four pillars)' -> 'thesis'."""
    cleaned = re.sub(r"\([^)]*\)", "", value).strip().lower()
    return cleaned


def _scan_slide(slide_id: str, body: str) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    body_lower = body.lower()
    for token, rule_id, label in (
        ("title", "SLIDE-TITLE", "Title"),
        ("primary job", "SLIDE-JOB", "Primary job"),
        ("core claim", "SLIDE-CLAIM", "Core claim"),
        ("support", "SLIDE-SUPPORT", "Support points"),
    ):
        if token not in body_lower:
            findings.append(("warn", rule_id, f"slide {slide_id}: missing '{label}' field"))

    if not any(token in body_lower for token in ("speaker notes", "presenter notes", "presenter intent")):
        findings.append(
            (
                "warn",
                "SLIDE-SPEAKER-NOTES",
                f"slide {slide_id}: missing presenter guidance ('Speaker notes' or 'Presenter intent')",
            )
        )

    title_value = _extract_field_value(body, "Title")
    if title_value and len(title_value) > TITLE_LENGTH_MAX:
        findings.append(
            (
                "warn",
                "SLIDE-TITLE-LENGTH",
                f"slide {slide_id}: title length {len(title_value)} exceeds {TITLE_LENGTH_MAX} chars",
            )
        )

    archetype_value = _extract_field_value(body, "Archetype")
    archetype_norm = _normalize_archetype(archetype_value)
    if archetype_value and archetype_norm not in KNOWN_ARCHETYPES:
        findings.append(
            (
                "warn",
                "SLIDE-ARCHETYPE-UNKNOWN",
                f"slide {slide_id}: archetype '{archetype_value}' is not in slide-archetypes.md",
            )
        )

    support_points = _extract_support_points(body)
    if support_points:
        bullets = support_points
        band_key = archetype_norm if archetype_norm in SUPPORT_COUNT_BANDS else None
        band = SUPPORT_COUNT_BANDS.get(band_key, SUPPORT_COUNT_DEFAULT) if band_key is not None else SUPPORT_COUNT_DEFAULT
        if band is None:
            return findings
        low, high = band
        if not (low <= len(bullets) <= high):
            findings.append(
                (
                    "warn",
                    "SLIDE-SUPPORT-COUNT",
                    f"slide {slide_id} ({archetype_norm or 'no-archetype'}): support points count {len(bullets)} outside {low}-{high} range",
                )
            )
    return findings


def _scan_requested_slide_count(target_slides: object, slide_count: int) -> list[tuple[str, str, str]]:
    if target_slides is None:
        return []
    if not isinstance(target_slides, int):
        return [("warn", "DECK-TARGET-SLIDES-TYPE", f"requested_output.target_slide_count should be an integer, got {type(target_slides).__name__}")]
    if target_slides <= 0:
        return [("warn", "DECK-TARGET-SLIDES-RANGE", f"requested_output.target_slide_count should be positive, got {target_slides}")]
    tolerance = 0 if target_slides == 1 else 1
    if abs(slide_count - target_slides) <= tolerance:
        return []
    return [
        (
            "warn",
            "DECK-TARGET-SLIDES",
            f"slide count {slide_count} is outside the requested target_slide_count {target_slides} ± {tolerance}",
        )
    ]


def _scan_slide_count(budget: str | None, slide_count: int) -> list[tuple[str, str, str]]:
    if not budget:
        return []
    band = SLIDE_COUNT_BANDS.get(budget)
    if band is None:
        return []
    low, high = band
    if low <= slide_count <= high:
        return []
    return [
        (
            "warn",
            "DECK-SLIDE-COUNT",
            f"slide count {slide_count} is outside the {budget} band ({low}-{high})",
        )
    ]


def _scan_prompts(run_dir: Path, storyboard_ids: list[str]) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    prompts_dir = run_dir / "prompts"
    if not prompts_dir.is_dir():
        findings.append(("error", "PROMPTS-001", "prompts/ directory missing"))
        return findings
    prompt_ids = sorted(p.stem for p in prompts_dir.glob("*.md") if p.name != "README.md")
    storyboard_set = set(storyboard_ids)
    for sid in storyboard_ids:
        if sid not in prompt_ids:
            findings.append(("error", "PROMPT-MISSING", f"slide {sid} has no prompt file at prompts/{sid}.md"))
    for pid in prompt_ids:
        if pid not in storyboard_set:
            findings.append(("warn", "PROMPT-ORPHAN", f"prompt prompts/{pid}.md has no matching storyboard slide"))
    return findings


def _scan_native_pptx_language(run_dir: Path) -> list[tuple[str, str, str]]:
    """Reject production artifacts that steer an image-first run back to native-PPTX."""
    findings: list[tuple[str, str, str]] = []
    candidate_files: list[Path] = []

    storyboard_path = run_dir / "work" / "storyboard.md"
    if storyboard_path.is_file():
        candidate_files.append(storyboard_path)

    prompts_dir = run_dir / "prompts"
    if prompts_dir.is_dir():
        candidate_files.extend(sorted(prompts_dir.glob("*.md")))

    for path in candidate_files:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue
        for line_number, line in enumerate(text.splitlines(), start=1):
            if NATIVE_PPTX_NEGATION_PATTERN.search(line):
                continue
            for pattern, reason in FORBIDDEN_PROMPT_PATTERNS:
                if not re.search(pattern, line, re.IGNORECASE):
                    continue
                findings.append(
                    (
                        "error",
                        "IMG-FIRST-NATIVE-PPTX-LANGUAGE",
                        f"{path.relative_to(run_dir)}:{line_number} {reason}; Deckit prompts/storyboards must stay image-first",
                    )
                )
                break
    return findings


def _scan_native_pptx_scripts(run_dir: Path) -> list[tuple[str, str, str]]:
    """Reject run-local scripts that look like native PowerPoint assembly."""
    findings: list[tuple[str, str, str]] = []
    scripts_dir = run_dir / "scripts"
    if not scripts_dir.is_dir():
        return findings

    script_suffixes = {".py", ".js", ".ts", ".mjs", ".cjs"}
    for script_path in sorted(p for p in scripts_dir.rglob("*") if p.is_file() and p.suffix.lower() in script_suffixes):
        try:
            text = script_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            continue

        matched_reasons = [
            reason
            for pattern, reason in NATIVE_PPTX_SCRIPT_PATTERNS
            if re.search(pattern, text, re.IGNORECASE | re.MULTILINE)
        ]
        if matched_reasons:
            findings.append(
                (
                    "error",
                    "IMG-FIRST-NATIVE-PPTX-SCRIPT",
                    f"{script_path.relative_to(run_dir)} {', '.join(matched_reasons)}; PPTX is only a container after $imagegen PNGs exist, never a native production route",
                )
            )
    return findings


def _audit_pptx_container(pptx_path: Path, label: str | None = None) -> tuple[list[tuple[str, str, str]], list[str]]:
    """Verify that a PPTX is an image-first container: exactly one picture per slide.

    Returns (findings, report_lines). Findings use the same (severity, rule_id, message)
    shape as review findings so the audit can be reused by `review`.
    """
    findings: list[tuple[str, str, str]] = []
    report_lines: list[str] = []
    display = label or str(pptx_path)

    if not pptx_path.is_file():
        return [("error", "PPTX-001", f"PPTX file not found: {display}")], [f"# PPTX Audit — {display}\n\nFile not found.\n"]

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX auditing.") from exc

    try:
        with zipfile.ZipFile(pptx_path) as package:
            bad_member = package.testzip()
            media = [name for name in package.namelist() if name.startswith("ppt/media/")]
    except zipfile.BadZipFile:
        return [("error", "PPTX-ZIP", f"{display} is not a valid ZIP/PPTX package")], [
            f"# PPTX Audit — {display}\n\nInvalid ZIP/PPTX package.\n"
        ]

    if bad_member is not None:
        findings.append(("error", "PPTX-ZIP-MEMBER", f"{display} failed ZIP integrity check at {bad_member}"))

    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        findings.append(("error", "PPTX-OPEN", f"{display} could not be opened by python-pptx: {exc}"))
        return findings, [f"# PPTX Audit — {display}\n\nCould not open with python-pptx: {exc}\n"]

    total_shapes = 0
    type_counts: dict[str, int] = {}
    slide_summaries: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        shape_types = [str(shape.shape_type) for shape in slide.shapes]
        total_shapes += len(shape_types)
        for shape_type in shape_types:
            type_counts[shape_type] = type_counts.get(shape_type, 0) + 1

        type_summary = ", ".join(f"{shape_type}: {count}" for shape_type, count in sorted({
            shape_type: shape_types.count(shape_type) for shape_type in set(shape_types)
        }.items()))
        slide_summaries.append(f"- Slide {idx}: {len(shape_types)} shape(s) — {type_summary or 'none'}\n")

        is_single_picture = len(slide.shapes) == 1 and shape_types and "PICTURE" in shape_types[0]
        if not is_single_picture:
            findings.append(
                (
                    "error",
                    "PPTX-NATIVE-CONTENT",
                    f"{display} slide {idx} has {len(shape_types)} shape(s) ({type_summary or 'none'}); expected exactly one full-slide picture and no native editable content",
                )
            )

    report_lines.append(f"# PPTX Audit — {display}\n\n")
    report_lines.append(f"- Slides: {len(prs.slides)}\n")
    report_lines.append(f"- Media assets: {len(media)}\n")
    report_lines.append(f"- Total slide shapes: {total_shapes}\n")
    report_lines.append(f"- Findings: {sum(1 for sev, _, _ in findings if sev == 'error')} error(s), {sum(1 for sev, _, _ in findings if sev == 'warn')} warning(s)\n\n")
    report_lines.append("## Shape Type Totals\n\n")
    if type_counts:
        for shape_type, count in sorted(type_counts.items()):
            report_lines.append(f"- {shape_type}: {count}\n")
    else:
        report_lines.append("- none\n")
    report_lines.append("\n## Slides\n\n")
    report_lines.extend(slide_summaries)
    if findings:
        report_lines.append("\n## Findings\n\n")
        for sev, rule_id, message in findings:
            report_lines.append(f"- **{sev.upper()}** `{rule_id}` — {message}\n")
    else:
        report_lines.append("\nAll checks passed: this PPTX is a one-image-per-slide container.\n")

    return findings, report_lines


def _find_pptx_deliverables(run_dir: Path) -> list[Path]:
    """Find likely PPTX deliverables in a run folder.

    Search both `dist/` and the run root because failed/manual runs often place
    deliverables at the root while still backfilling Deckit-like artifacts later.
    Ignore PowerPoint lock files (`~$*.pptx`).
    """
    candidates: list[Path] = []
    for folder in (run_dir / "dist", run_dir):
        if folder.is_dir():
            candidates.extend(p for p in folder.glob("*.pptx") if not p.name.startswith("~$"))
    return sorted(set(candidates))


def _scan_image_first_artifacts(run_dir: Path, storyboard_ids: list[str]) -> list[tuple[str, str, str]]:
    findings: list[tuple[str, str, str]] = []
    generated_dir = run_dir / "assets" / "generated-slides"
    generated_pngs = sorted(generated_dir.glob("*.png")) if generated_dir.is_dir() else []
    generated_ids = {p.stem for p in generated_pngs}
    dist_pptx = _find_pptx_deliverables(run_dir)
    preview_path = run_dir / "dist" / "preview.png"

    if dist_pptx and not generated_pngs:
        pptx_names = ", ".join(p.name for p in dist_pptx)
        findings.append(
            (
                "error",
                "IMG-FIRST-PPTX-WITHOUT-GENERATED-SLIDES",
                f"dist contains PPTX deliverable(s) ({pptx_names}) but assets/generated-slides has no PNGs from $imagegen",
            )
        )
    elif dist_pptx and storyboard_ids and set(storyboard_ids) - generated_ids:
        missing = ", ".join(sorted(set(storyboard_ids) - generated_ids))
        pptx_names = ", ".join(p.name for p in dist_pptx)
        findings.append(
            (
                "error",
                "IMG-FIRST-PPTX-BEFORE-COMPLETE-IMAGES",
                f"dist contains PPTX deliverable(s) ({pptx_names}) before every storyboard slide has a generated PNG; missing: {missing}",
            )
        )

    for pptx_path in dist_pptx:
        audit_findings, _ = _audit_pptx_container(pptx_path, str(pptx_path.relative_to(run_dir)))
        findings.extend(audit_findings)

    for sid in storyboard_ids:
        if generated_pngs and sid not in generated_ids:
            findings.append(
                (
                    "error",
                    "IMG-FIRST-GENERATED-SLIDE-MISSING",
                    f"slide {sid} has no generated image at assets/generated-slides/{sid}.png",
                )
            )
    for image_id in sorted(generated_ids - set(storyboard_ids)):
        findings.append(
            (
                "warn",
                "IMG-FIRST-GENERATED-SLIDE-ORPHAN",
                f"generated image assets/generated-slides/{image_id}.png has no matching storyboard slide",
            )
        )

    if storyboard_ids and set(storyboard_ids).issubset(generated_ids) and not preview_path.is_file():
        findings.append(
            (
                "warn",
                "IMG-FIRST-PREVIEW-MISSING",
                "generated slide PNGs are complete but the standard preview is missing at dist/preview.png",
            )
        )

    findings.extend(_scan_native_pptx_scripts(run_dir))
    findings.extend(_scan_native_pptx_language(run_dir))
    return findings


def _package_images(run_dir: Path, out_path: Path | None) -> int:
    """Package generated slide PNGs into a non-editable image-only PPTX container."""
    run_dir = run_dir.resolve()
    if not run_dir.is_dir():
        raise NotADirectoryError(f"run directory does not exist: {run_dir}")

    slide_blocks, sb_findings = _parse_storyboard(run_dir / "work" / "storyboard.md")
    errors = [finding for finding in sb_findings if finding[0] == "error"]
    if errors:
        detail = "; ".join(message for _, _, message in errors)
        raise ValueError(f"cannot package images because storyboard is invalid: {detail}")

    slide_ids = [slide_id for slide_id, _ in slide_blocks]
    if not slide_ids:
        raise ValueError("cannot package images because storyboard contains no slides")

    generated_dir = run_dir / "assets" / "generated-slides"
    missing = [slide_id for slide_id in slide_ids if not (generated_dir / f"{slide_id}.png").is_file()]
    if missing:
        raise FileNotFoundError(
            "cannot package PPTX before every storyboard slide has a generated PNG. Missing: "
            + ", ".join(f"assets/generated-slides/{slide_id}.png" for slide_id in missing)
        )

    if out_path is None:
        out_path = run_dir / "dist" / f"{run_dir.name}.pptx"
    elif not out_path.is_absolute():
        out_path = (run_dir / out_path).resolve()
    else:
        out_path = out_path.resolve()

    dist_dir = run_dir / "dist"
    dist_dir.mkdir(parents=True, exist_ok=True)
    _ensure_inside(dist_dir, out_path)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for stable PPTX container packaging. "
            "Install the dev tools with their dependencies, then rerun deckit-dev package-images."
        ) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    notes_by_slide_id = {slide_id: _build_speaker_notes(slide_id, body) for slide_id, body in slide_blocks}

    for slide_id in slide_ids:
        slide = prs.slides.add_slide(blank_layout)
        image_path = generated_dir / f"{slide_id}.png"
        # Packaging discipline: one generated full-slide image per slide; no editable text boxes/charts/native layouts.
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        # Add PowerPoint speaker notes from the storyboard so presenters have a talk track
        # without changing the image-first, non-editable slide canvas.
        slide.notes_slide.notes_text_frame.text = notes_by_slide_id[slide_id]

    prs.save(str(out_path))

    audit_findings, _ = _audit_pptx_container(out_path)
    audit_errors = [message for sev, _, message in audit_findings if sev == "error"]
    if audit_errors:
        raise ValueError("packaged PPTX failed image-first container audit: " + "; ".join(audit_errors))

    check = Presentation(str(out_path))
    if len(check.slides) != len(slide_ids):
        raise ValueError(f"packaged PPTX has {len(check.slides)} slide(s), expected {len(slide_ids)}")
    for idx, slide in enumerate(check.slides, start=1):
        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            raise ValueError(f"packaged PPTX slide {idx} is missing speaker notes")

    preview_path = _package_preview(run_dir, None)

    print(f"wrote: {out_path}")
    print(f"slides packaged: {len(slide_ids)}")
    print("mode: image-first PNG container; no editable native PowerPoint slide content; speaker notes included")
    print(f"preview: {preview_path}")
    return 0


def _find_pdf_deliverables(run_dir: Path) -> list[Path]:
    candidates: list[Path] = []
    for folder in (run_dir / "dist", run_dir):
        if folder.is_dir():
            candidates.extend(p for p in folder.glob("*.pdf") if not p.name.startswith("~$"))
    return sorted(set(candidates))


def _scan_requested_delivery_artifact(run_dir: Path, delivery_target: object) -> list[tuple[str, str, str]]:
    if delivery_target is None:
        return []
    if delivery_target == "pptx":
        if not _find_pptx_deliverables(run_dir):
            return [("error", "DELIVERY-PPTX-MISSING", "requested final delivery target is pptx but no PPTX was found in dist/")]
        return []
    if delivery_target == "pdf":
        if not _find_pdf_deliverables(run_dir):
            return [("error", "DELIVERY-PDF-MISSING", "requested final delivery target is pdf but no PDF was found in dist/")]
        return []
    return []


def _require_generated_slide_images(run_dir: Path) -> tuple[list[str], list[Path]]:
    """Return storyboard-ordered generated slide PNGs, failing if any are missing."""
    run_dir = run_dir.resolve()
    if not run_dir.is_dir():
        raise NotADirectoryError(f"run directory does not exist: {run_dir}")

    slide_blocks, sb_findings = _parse_storyboard(run_dir / "work" / "storyboard.md")
    errors = [finding for finding in sb_findings if finding[0] == "error"]
    if errors:
        detail = "; ".join(message for _, _, message in errors)
        raise ValueError(f"cannot package generated slide images because storyboard is invalid: {detail}")

    slide_ids = [slide_id for slide_id, _ in slide_blocks]
    if not slide_ids:
        raise ValueError("cannot package generated slide images because storyboard contains no slides")

    generated_dir = run_dir / "assets" / "generated-slides"
    image_paths = [generated_dir / f"{slide_id}.png" for slide_id in slide_ids]
    missing = [path for path in image_paths if not path.is_file()]
    if missing:
        raise FileNotFoundError(
            "cannot package before every storyboard slide has a generated PNG. Missing: "
            + ", ".join(str(path.relative_to(run_dir)) for path in missing)
        )
    return slide_ids, image_paths


def _resolve_dist_output(run_dir: Path, out_path: Path | None, suffix: str) -> Path:
    if out_path is None:
        out_path = run_dir / "dist" / f"{run_dir.name}{suffix}"
    elif not out_path.is_absolute():
        out_path = (run_dir / out_path).resolve()
    else:
        out_path = out_path.resolve()

    dist_dir = run_dir / "dist"
    dist_dir.mkdir(parents=True, exist_ok=True)
    _ensure_inside(dist_dir, out_path)
    return out_path


def _package_pdf(run_dir: Path, out_path: Path | None) -> int:
    """Package generated slide PNGs into an image-only PDF, one slide image per page."""
    run_dir = run_dir.resolve()
    slide_ids, image_paths = _require_generated_slide_images(run_dir)
    out_path = _resolve_dist_output(run_dir, out_path, ".pdf")

    try:
        from PIL import Image, JpegImagePlugin  # noqa: F401
    except ImportError as exc:
        raise RuntimeError("Pillow is required for PDF packaging.") from exc

    rgb_pages = []
    for image_path in image_paths:
        with Image.open(image_path) as image:
            page = image.convert("RGB")
            rgb_pages.append(page.copy())

    first, rest = rgb_pages[0], rgb_pages[1:]
    first.save(out_path, save_all=True, append_images=rest)

    preview_path = _package_preview(run_dir, None)

    print(f"wrote: {out_path}")
    print(f"pages packaged: {len(slide_ids)}")
    print("mode: image-first PNG pages packaged as an image-only PDF")
    print(f"preview: {preview_path}")
    return 0


def _package_preview(run_dir: Path, out_path: Path | None, width: int = 900, gap: int = 24, background: str = "white") -> Path:
    """Create the standard medium-size vertical preview PNG from generated slide images."""
    run_dir = run_dir.resolve()
    slide_ids, image_paths = _require_generated_slide_images(run_dir)
    if out_path is None:
        out_path = run_dir / "dist" / "preview.png"
    else:
        out_path = _resolve_dist_output(run_dir, out_path, ".png")
    out_path = out_path.resolve()
    dist_dir = run_dir / "dist"
    dist_dir.mkdir(parents=True, exist_ok=True)
    _ensure_inside(dist_dir, out_path)
    if width <= 0:
        raise ValueError(f"width must be positive, got {width}")
    if gap < 0:
        raise ValueError(f"gap must be non-negative, got {gap}")

    try:
        from PIL import Image, ImageColor
    except ImportError as exc:
        raise RuntimeError("Pillow is required for preview image packaging.") from exc

    try:
        fill = ImageColor.getrgb(background)
    except ValueError as exc:
        raise ValueError(f"background must be a Pillow-compatible color, got {background!r}") from exc

    preview_slides = []
    for image_path in image_paths:
        with Image.open(image_path) as image:
            slide = image.convert("RGB")
            if slide.width > width:
                height = round(slide.height * (width / slide.width))
                slide = slide.resize((width, height), Image.Resampling.LANCZOS)
            preview_slides.append(slide.copy())

    max_width = max(image.width for image in preview_slides)
    total_height = sum(image.height for image in preview_slides) + gap * (len(preview_slides) - 1)
    canvas = Image.new("RGB", (max_width, total_height), fill)

    y = 0
    for image in preview_slides:
        x = (max_width - image.width) // 2
        canvas.paste(image, (x, y))
        y += image.height + gap

    canvas.save(out_path)
    return out_path


def _package_preview_command(run_dir: Path, out_path: Path | None, width: int, gap: int, background: str) -> int:
    """CLI wrapper for the standard vertical preview PNG."""
    preview_path = _package_preview(run_dir, out_path, width=width, gap=gap, background=background)
    print(f"wrote: {preview_path}")
    print(f"preview: {preview_path}")
    print("mode: standard medium-size vertical preview generated from slide PNGs")
    return 0


def audit_pptx(pptx_path: Path) -> int:
    pptx_path = pptx_path.resolve()
    findings, report_lines = _audit_pptx_container(pptx_path)
    print("".join(report_lines), end="")
    counts = {"error": 0, "warn": 0}
    for sev, _, _ in findings:
        counts[sev] = counts.get(sev, 0) + 1
    return 0 if counts["error"] == 0 else 1


def review(run_dir: Path) -> int:
    run_dir = run_dir.resolve()
    if not run_dir.is_dir():
        raise NotADirectoryError(f"run directory does not exist: {run_dir}")

    findings: list[tuple[str, str, str]] = []

    run_json_path = run_dir / "run.json"
    manifest: dict[str, object] = {}
    if not run_json_path.is_file():
        findings.append(("error", "RUN-001", "run.json missing"))
    else:
        try:
            manifest = json.loads(run_json_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            findings.append(("error", "RUN-002", f"run.json is not valid JSON: {exc}"))

    mode = manifest.get("production_mode") if isinstance(manifest, dict) else None
    if mode is None:
        findings.append(("warn", "MODE-001", "production_mode is null in run.json"))
    elif mode not in PRODUCTION_MODES:
        findings.append(("error", "MODE-002", f"production_mode '{mode}' is not one of {PRODUCTION_MODES}"))

    findings.extend(_scan_brief(run_dir / "work" / "deck-brief.md"))

    slide_blocks, sb_findings = _parse_storyboard(run_dir / "work" / "storyboard.md")
    findings.extend(sb_findings)
    for slide_id, body in slide_blocks:
        findings.extend(_scan_slide(slide_id, body))

    storyboard_ids = [sid for sid, _ in slide_blocks]
    budget = manifest.get("budget_mode") if isinstance(manifest, dict) else None
    requested_output = manifest.get("requested_output") if isinstance(manifest, dict) else None
    delivery_target = requested_output.get("delivery_target") if isinstance(requested_output, dict) else None
    target_slides = requested_output.get("target_slide_count") if isinstance(requested_output, dict) else None
    if delivery_target is not None and delivery_target not in DELIVERY_TARGETS:
        findings.append(
            (
                "warn",
                "DECK-DELIVERY-TARGET",
                f"requested_output.delivery_target '{delivery_target}' is not one of {DELIVERY_TARGETS}",
            )
        )
    else:
        findings.extend(_scan_requested_delivery_artifact(run_dir, delivery_target))
    if storyboard_ids:
        findings.extend(_scan_requested_slide_count(target_slides, len(storyboard_ids)))
    if storyboard_ids and isinstance(budget, str) and target_slides is None:
        findings.extend(_scan_slide_count(budget, len(storyboard_ids)))
    findings.extend(_scan_prompts(run_dir, storyboard_ids))
    if mode == "image-first":
        findings.extend(_scan_image_first_artifacts(run_dir, storyboard_ids))

    counts = {"error": 0, "warn": 0}
    for sev, _, _ in findings:
        counts[sev] = counts.get(sev, 0) + 1

    dist_dir = run_dir / "dist"
    dist_dir.mkdir(parents=True, exist_ok=True)
    review_path = dist_dir / "review.md"

    lines: list[str] = []
    lines.append(f"# Review Report — {run_dir.name}\n\n")
    lines.append(f"- Production mode: `{mode if mode else 'unset'}`\n")
    if isinstance(requested_output, dict):
        lines.append(f"- Final delivery target: `{requested_output.get('delivery_target') or 'unset'}`\n")
        lines.append(f"- Target slide count: `{requested_output.get('target_slide_count') or 'unset'}`\n")
    lines.append(f"- Slides found in storyboard: {len(storyboard_ids)}\n")
    lines.append(f"- Findings: {counts['error']} error(s), {counts['warn']} warning(s)\n\n")
    if not findings:
        lines.append("All checks passed.\n")
    else:
        lines.append("## Findings\n\n")
        for sev, rule_id, msg in findings:
            lines.append(f"- **{sev.upper()}** `{rule_id}` — {msg}\n")
    review_path.write_text("".join(lines), encoding="utf-8")

    print(f"wrote: {review_path}")
    print(f"errors: {counts['error']}; warnings: {counts['warn']}")
    return 0 if counts["error"] == 0 else 1


def main() -> int:
    parser = argparse.ArgumentParser(description="Deckit development tools")
    subparsers = parser.add_subparsers(dest="command", required=True)

    inspect_parser = subparsers.add_parser("inspect", help="Inspect the plugin manifest and skills")
    inspect_parser.add_argument("--plugin", type=Path, default=DEFAULT_PLUGIN)

    marketplace_parser = subparsers.add_parser("inspect-marketplace", help="Inspect repo-local plugin marketplace")
    marketplace_parser.add_argument("--marketplace", type=Path, default=DEFAULT_MARKETPLACE)

    new_run_parser = subparsers.add_parser("new-run", help="Create a standard local run directory from a text, PDF, or URL source")
    new_run_parser.add_argument("--source", required=True, help="Path to a text/Markdown file, a .pdf file, or an http(s) URL")
    new_run_parser.add_argument("--name", help="Run name. Defaults to the source filename stem (or URL host+path).")
    new_run_parser.add_argument(
        "--runs-dir",
        type=Path,
        default=DEFAULT_RUNS_DIR,
        help="Directory that will contain run folders. Defaults to ./outputs in the current working directory; pass the caller workspace's outputs directory when invoking from a plugin/dev checkout.",
    )
    new_run_parser.add_argument("--force", action="store_true", help="Reuse an existing run directory and replace source copy")
    new_run_parser.add_argument(
        "--mode",
        choices=PRODUCTION_MODES,
        default="image-first",
        help="Deprecated compatibility flag. Deckit always uses image-first; no other production modes are valid.",
    )
    new_run_parser.add_argument("--budget", choices=("quick", "balanced", "premium"), help="Budget mode (recorded in run.json).")
    new_run_parser.add_argument("--delivery-target", choices=DELIVERY_TARGETS, help="Requested final delivery target recorded in run.json.")
    new_run_parser.add_argument("--target-slides", type=int, help="Approximate target slide/page count recorded in run.json.")

    review_parser = subparsers.add_parser("review", help="Run rule-based quality checks on a run folder and write dist/review.md")
    review_parser.add_argument("--run", type=Path, required=True, help="Path to the run directory to review.")

    package_parser = subparsers.add_parser(
        "package-images",
        help="Package assets/generated-slides/*.png into a non-editable image-only PPTX container",
    )
    package_parser.add_argument("--run", type=Path, required=True, help="Path to the run directory to package.")
    package_parser.add_argument(
        "--out",
        type=Path,
        help="Output PPTX path. Relative paths are resolved inside the run directory; default: dist/<run-name>.pptx.",
    )

    pdf_parser = subparsers.add_parser(
        "package-pdf",
        help="Package assets/generated-slides/*.png into an image-only PDF, one slide image per page",
    )
    pdf_parser.add_argument("--run", type=Path, required=True, help="Path to the run directory to package.")
    pdf_parser.add_argument(
        "--out",
        type=Path,
        help="Output PDF path. Relative paths are resolved inside the run directory; default: dist/<run-name>.pdf.",
    )

    preview_parser = subparsers.add_parser(
        "package-preview",
        help="Create the standard medium-size vertical preview PNG from assets/generated-slides/*.png",
    )
    preview_parser.add_argument("--run", type=Path, required=True, help="Path to the run directory to package.")
    preview_parser.add_argument(
        "--out",
        type=Path,
        help="Output PNG path. Relative paths are resolved inside the run directory; default: dist/preview.png.",
    )
    preview_parser.add_argument("--width", type=int, default=900, help="Maximum preview width in pixels.")
    preview_parser.add_argument("--gap", type=int, default=24, help="Vertical gap in pixels between slides.")
    preview_parser.add_argument("--background", default="white", help="Canvas background color for gaps and side padding.")

    audit_parser = subparsers.add_parser(
        "audit-pptx",
        help="Audit a PPTX for Deckit image-first compliance (exactly one full-slide picture per slide)",
    )
    audit_parser.add_argument("--pptx", type=Path, required=True, help="PPTX file to audit.")

    ingest_parser = subparsers.add_parser("ingest", help="Convert a PDF or URL into source/input.md (Markdown)")
    ingest_group = ingest_parser.add_mutually_exclusive_group(required=True)
    ingest_group.add_argument("--pdf", type=Path, help="Path to a .pdf file")
    ingest_group.add_argument("--url", help="An http(s) URL")
    ingest_parser.add_argument("--out", type=Path, required=True, help="Output Markdown path (e.g. source/input.md)")

    args = parser.parse_args()

    if args.command == "inspect":
        return inspect_plugin(args.plugin.resolve())
    if args.command == "inspect-marketplace":
        return inspect_marketplace(args.marketplace)
    if args.command == "new-run":
        return new_run(
            args.source,
            args.name,
            args.runs_dir,
            args.force,
            args.mode,
            args.budget,
            args.delivery_target,
            args.target_slides,
        )
    if args.command == "review":
        return review(args.run)
    if args.command == "package-images":
        return _package_images(args.run, args.out)
    if args.command == "package-pdf":
        return _package_pdf(args.run, args.out)
    if args.command == "package-preview":
        return _package_preview_command(args.run, args.out, args.width, args.gap, args.background)
    if args.command == "audit-pptx":
        return audit_pptx(args.pptx)
    if args.command == "ingest":
        from deckit_dev.ingest import ingest_pdf, ingest_url

        if args.pdf is not None:
            return ingest_pdf(args.pdf, args.out)
        return ingest_url(args.url, args.out)

    parser.error(f"unknown command: {args.command}")
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
